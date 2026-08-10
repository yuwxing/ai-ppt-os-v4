# -*- coding: utf-8 -*-
"""
build_pptx.py — 初中英语话题课件生成器（母版占位符版）

设计要点
--------
1. 不再手工画框：所有页面都从 template.pptx 的自定义版式派生，
   版面、字号、字色、留白由母版统一控制，因此每页天然对齐、风格一致。
2. 图片 / 音视频一律使用【真实占位符】：
     PICTURE(18)    → 教师在 PowerPoint 中点击图标，弹出【本地文件选择器】
     MEDIA_CLIP(10) → 点击弹出【本地音频/视频选择器】
   不再写 file:/// 超链接（那会被 PowerPoint 当成网页链接）。
   未插入素材的占位符在放映和打印时自动隐藏，不会留下空白框。
3. 练习页答案框命名为 REVEAL_ANSWER，由 postprocess_anim.ps1 加"单击才出现"动画。

用法
----
python build_pptx.py --json content.json --out 课件.pptx
"""

import argparse
import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_TEMPLATE = os.path.join(os.path.dirname(HERE), "template.pptx")

FONT = "Microsoft YaHei"
FONT_EN = "Segoe UI"

NAVY = RGBColor(0x12, 0x39, 0x5E)
DARK = RGBColor(0x16, 0x18, 0x1A)
GRAY = RGBColor(0x5A, 0x64, 0x70)
GREEN = RGBColor(0x0F, 0x6B, 0x34)
ACCENT = RGBColor(0xE2, 0x60, 0x1A)
PURPLE = RGBColor(0x6B, 0x3F, 0xA0)

LINE_SPACING = 1.32

# 需要教师自备素材的清单（写入 README）
ASSET_NOTES = []


# ---------------------------------------------------------------- 基础工具
def layout_by_name(prs, name):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    raise KeyError(f"版式不存在: {name}（模板: {prs}）")


def ph_by_idx(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def ph_by_type(slide, pptype, nth=0):
    hits = [p for p in slide.placeholders if p.placeholder_format.type == pptype]
    return hits[nth] if len(hits) > nth else None


def drop(shape):
    """删除未使用的占位符，避免编辑视图里出现多余提示框。"""
    if shape is not None:
        shape._element.getparent().remove(shape._element)


def _apply_font(run, size, color, bold, italic=False, font=FONT):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font
    # 同时设置东亚字体，避免中文回退到宋体
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def _no_bullet(para):
    pPr = para._p.get_or_add_pPr()
    if pPr.find(qn("a:buNone")) is None:
        bu = pPr.makeelement(qn("a:buNone"), {})
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is not None:
            defRPr.addprevious(bu)
        else:
            pPr.append(bu)


def write_lines(ph, lines, align=PP_ALIGN.LEFT):
    """
    lines: [{text, size, color, bold, italic, space_after, space_before, indent, font}]
    """
    tf = ph.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = ln.get("align", align)
        para.line_spacing = ln.get("line_spacing", LINE_SPACING)
        if ln.get("space_before"):
            para.space_before = Pt(ln["space_before"])
        para.space_after = Pt(ln.get("space_after", 8))
        if ln.get("indent"):
            para.level = ln["indent"]
        _no_bullet(para)
        run = para.add_run()
        run.text = ln["text"]
        _apply_font(
            run,
            ln.get("size", 24),
            ln.get("color", DARK),
            ln.get("bold", False),
            ln.get("italic", False),
            ln.get("font", FONT),
        )


def set_title(slide, text, size=None, color=NAVY):
    t = ph_by_type(slide, PP_PLACEHOLDER.TITLE) or ph_by_type(slide, PP_PLACEHOLDER.CENTER_TITLE)
    if t is None:
        return
    tf = t.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = tf.paragraphs[0].alignment
    _no_bullet(p)
    run = p.add_run()
    run.text = text
    _apply_font(run, size or 38, color, True)


def add_textbox(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, name=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if name:
        tb.name = name
    write_lines(tb, lines, align=align)
    return tb


def note_asset(page, what):
    ASSET_NOTES.append(f"第 {page:>2} 页 | {what}")


def keep_picture(slide, page, what, nth=0):
    """保留图片占位符（教师点击插入本地图片），并登记到素材清单。"""
    p = ph_by_type(slide, PP_PLACEHOLDER.PICTURE, nth)
    if p is not None:
        note_asset(page, what)
    return p


def keep_media(slide, page, what):
    p = ph_by_type(slide, PP_PLACEHOLDER.MEDIA_CLIP)
    if p is not None:
        note_asset(page, what)
    return p


# ---------------------------------------------------------------- 页面构建
def build(data, prs):
    _T = data.get("titles", {})
    def T(key, default):
        return _T.get(key, default)
    meta = data.get("meta", {})
    author = meta.get("author", "")
    topic_en = meta.get("topic_en", "")
    topic_cn = meta.get("topic_cn", "")

    L = lambda n: layout_by_name(prs, n)
    page = 0

    def new(layout_name):
        nonlocal page
        page += 1
        return prs.slides.add_slide(L(layout_name)), page

    # ---- 1 封面 ----
    s, pg = new("L_Cover")
    set_title(s, meta.get("cover_title", topic_en), size=54)
    sub = ph_by_idx(s, 10)
    write_lines(sub, [{"text": meta.get("cover_subtitle", topic_cn),
                       "size": 26, "color": GRAY, "align": PP_ALIGN.CENTER}],
                align=PP_ALIGN.CENTER)
    keep_picture(s, pg, T("cover_pic", "封面主视觉图（建议：与话题相关的高清横构图）"))
    if author:
        add_textbox(s, 9.6, 6.72, 3.1, 0.5,
                    [{"text": f"授课教师：{author}", "size": 15, "color": GRAY,
                      "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)

    # ---- 2 学习目标 ----
    s, pg = new("L_TitleBody")
    set_title(s, T("objectives", "学习目标  Learning Objectives"))
    body = ph_by_idx(s, 10)
    lines = []
    for i, o in enumerate(data.get("objectives", []), 1):
        lines.append({"text": f"{i}.  {o}", "size": 25, "space_after": 14})
    write_lines(body, lines)

    # ---- 3 话题导入 ----
    s, pg = new("L_TitleImage")
    set_title(s, T("lead_in", "话题导入  Lead-in"))
    body = ph_by_idx(s, 10)
    li = data.get("lead_in", {})
    lines = []
    for q in li.get("questions", []):
        lines.append({"text": f"?  {q}", "size": 26, "color": NAVY,
                      "bold": True, "space_after": 14})
    if li.get("intro"):
        lines.append({"text": li["intro"], "size": 21, "color": GRAY,
                      "space_before": 12, "space_after": 6})
    write_lines(body, lines)
    keep_picture(s, pg, T("lead_in_pic", "导入情境图（建议：与话题相关的场景图）"))

    # ---- 4 Task 1 词汇卡 ----
    s, pg = new("L_VocabCards")
    t1 = data.get("task1", {})
    set_title(s, f"{t1.get('label','Task 1')}  ·  {t1.get('desc','')}")
    xs = [0.70, 4.67, 8.64]
    ys_cap = [3.42, 6.00]
    items = t1.get("items", [])[:6]
    for i, it in enumerate(items):
        col, row = i % 3, i // 3
        add_textbox(s, xs[col], ys_cap[row], 3.99, 0.80,
                    [{"text": it.get("word", ""), "size": 23, "color": NAVY,
                      "bold": True, "align": PP_ALIGN.CENTER, "space_after": 2,
                      "font": FONT_EN},
                     {"text": it.get("cn", ""), "size": 18, "color": GRAY,
                      "align": PP_ALIGN.CENTER, "space_after": 0}],
                    align=PP_ALIGN.CENTER)
    # 多余的图片占位符删掉
    pics = [p for p in s.placeholders if p.placeholder_format.type == PP_PLACEHOLDER.PICTURE]
    for i, p in enumerate(pics):
        if i < len(items):
            note_asset(pg, f"词汇配图 {i+1}：{items[i].get('word','')}")
        else:
            drop(p)

    # ---- 5 Brain storm ----
    s, pg = new("L_TwoCol")
    bs = data.get("brainstorm", {})
    set_title(s, bs.get("title", "Brain storm"))
    items = bs.get("items", [])
    half = (len(items) + 1) // 2
    for idx, chunk in ((10, items[:half]), (11, items[half:])):
        lines = []
        for it in chunk:
            lines.append({"text": it.get("word", ""), "size": 26, "color": NAVY,
                          "bold": True, "space_after": 2, "font": FONT_EN})
            lines.append({"text": it.get("cn", ""), "size": 19, "color": GRAY,
                          "space_after": 16})
        write_lines(ph_by_idx(s, idx), lines)

    # ---- 6 Let's read（点击出答案）----
    s, pg = new("L_Exercise")
    lr = data.get("lets_read", {})
    set_title(s, lr.get("title", "Let's read"))
    q_lines, a_lines = [], []
    for i, it in enumerate(lr.get("items", []), 1):
        q_lines.append({"text": f"{i}.  {it.get('en','')}", "size": 26,
                        "space_after": 14, "font": FONT_EN})
        a_lines.append({"text": f"{i}. {it.get('cn','')}", "size": 20,
                        "color": GREEN, "bold": True, "space_after": 6})
    write_lines(ph_by_idx(s, 10), q_lines)
    ans = ph_by_idx(s, 11)
    write_lines(ans, [{"text": "答案 / 解析", "size": 20, "color": GREEN,
                       "bold": True, "space_after": 8}] + a_lines)
    ans.name = "REVEAL_ANSWER"

    # ---- 7 Culture Link ----
    s, pg = new("L_TitleImage")
    set_title(s, T("culture", "Culture Link  ·  文化链接"))
    cu = data.get("culture", {})
    write_lines(ph_by_idx(s, 10), [
        {"text": cu.get("question", ""), "size": 26, "color": NAVY,
         "bold": True, "space_after": 18},
        {"text": cu.get("fact", ""), "size": 23, "space_after": 6},
    ])
    keep_picture(s, pg, T("culture_pic", "文化链接配图（建议：与话题相关的真实图片）"))

    # ---- 8 Task 2 对话补全（点击出答案）----
    s, pg = new("L_Exercise")
    t2 = data.get("task2", {})
    set_title(s, f"{t2.get('label','Task 2')}  ·  {t2.get('tag','')}", color=ACCENT)
    q_lines, a_lines = [], []
    for i, ln in enumerate(t2.get("lines", []), 1):
        q_lines.append({"text": f"{ln.get('speaker','')}  {ln.get('text','')}",
                        "size": 26, "space_after": 12, "font": FONT_EN})
        a_lines.append({"text": f"{i}. {ln.get('answer','')}", "size": 20,
                        "color": GREEN, "bold": True, "space_after": 4})
    write_lines(ph_by_idx(s, 10), q_lines)
    ans = ph_by_idx(s, 11)
    tips = t2.get("tips", "")
    extra = [{"text": f"Tips：{tips}", "size": 17, "color": GRAY,
              "bold": False, "space_before": 6, "space_after": 0}] if tips else []
    write_lines(ans, [{"text": "答案", "size": 20, "color": GREEN,
                       "bold": True, "space_after": 6}] + a_lines + extra)
    ans.name = "REVEAL_ANSWER"

    # ---- 9 功能表达 ----
    s, pg = new("L_TwoCol")
    set_title(s, T("expressions", "功能表达  Useful Expressions"))
    ex = data.get("expressions", {})
    left = [{"text": "问  Asking", "size": 24, "color": ACCENT, "bold": True,
             "space_after": 14}]
    for e in ex.get("ask", []):
        left.append({"text": f"·  {e}", "size": 24, "space_after": 14,
                     "font": FONT_EN})
    right = [{"text": "答  Answering", "size": 24, "color": ACCENT, "bold": True,
              "space_after": 14}]
    for e in ex.get("answer", []):
        right.append({"text": f"·  {e}", "size": 24, "space_after": 14,
                      "font": FONT_EN})
    write_lines(ph_by_idx(s, 10), left)
    write_lines(ph_by_idx(s, 11), right)

    # ---- 10 Task 3 项目式任务 ----
    s, pg = new("L_TitleBody")
    t3 = data.get("task3", {})
    set_title(s, f"{t3.get('label','Task 3')}  ·  {t3.get('tag','')}", color=PURPLE)
    lines = []
    if t3.get("time_limit"):
        lines.append({"text": t3["time_limit"], "size": 20, "color": ACCENT,
                      "bold": True, "space_after": 16})
    for sc in t3.get("scenes", []):
        lines.append({"text": f"·  {sc}", "size": 24, "space_after": 18})
    write_lines(ph_by_idx(s, 10), lines)

    # ---- 11 直击中考 · 听力（媒体占位符）----
    lst = data.get("listening", [])
    item = lst[0] if lst else {}
    s, pg = new("L_Media")
    set_title(s, T("media", data.get("listening_intro", "直击中考 · 听力实战")))
    lines = [{"text": item.get("title", ""), "size": 21, "color": ACCENT,
              "bold": True, "space_after": 14}]
    if item.get("stem"):
        lines.append({"text": item["stem"], "size": 24, "space_after": 14})
    if item.get("question"):
        lines.append({"text": item["question"], "size": 26, "color": NAVY,
                      "bold": True, "space_after": 12, "font": FONT_EN})
    for op in item.get("options", []):
        lines.append({"text": f"    {op}", "size": 24, "space_after": 8,
                      "font": FONT_EN})
    if item.get("source"):
        lines.append({"text": item["source"], "size": 15, "color": GRAY,
                      "space_before": 10, "space_after": 0})
    write_lines(ph_by_idx(s, 10), lines)
    keep_media(s, pg, T("media_note", "音频/视频（mp3/wav/mp4，点击媒体占位符插入本地文件）"))

    # ---- 12 听力解析（点击出答案）----
    s, pg = new("L_Exercise")
    set_title(s, T("analysis", "听力解析  Script & Key"))
    write_lines(ph_by_idx(s, 10), [
        {"text": "听力原文 Script", "size": 20, "color": ACCENT, "bold": True,
         "space_after": 10},
        {"text": item.get("script", ""), "size": 24, "space_after": 6,
         "font": FONT_EN},
    ])
    ans = ph_by_idx(s, 11)
    write_lines(ans, [
        {"text": "答题技巧", "size": 20, "color": GREEN, "bold": True,
         "space_after": 8},
        {"text": item.get("tip", ""), "size": 22, "color": GREEN, "space_after": 0},
    ])
    ans.name = "REVEAL_ANSWER"

    # ---- 13 技巧秘籍 ----
    s, pg = new("L_TitleBody")
    ts = data.get("tips_summary", {})
    set_title(s, ts.get("title", "秘籍只招"))
    lines = [{"text": f"·  {x}", "size": 24, "space_after": 18}
             for x in ts.get("items", [])]
    write_lines(ph_by_idx(s, 10), lines)

    # ---- 14 作业 ----
    s, pg = new("L_TitleBody")
    set_title(s, T("homework", "课后作业  Homework"))
    lines = [{"text": f"{i}.  {x}", "size": 25, "space_after": 20}
             for i, x in enumerate(data.get("homework", {}).get("items", []), 1)]
    write_lines(ph_by_idx(s, 10), lines)

    # ---- 15 情感升华 ----
    s, pg = new("L_TitleOnly")
    cl = data.get("closing", {}).get("lines", [])
    set_title(s, cl[0] if cl else topic_en, size=46)
    write_lines(ph_by_idx(s, 10),
                [{"text": cl[1] if len(cl) > 1 else "", "size": 30,
                  "color": ACCENT, "bold": True, "align": PP_ALIGN.CENTER}],
                align=PP_ALIGN.CENTER)

    # ---- 16 致谢 ----
    s, pg = new("L_TitleOnly")
    set_title(s, data.get("thanks", "Thanks for listening!"), size=46)
    sub = ph_by_idx(s, 10)
    write_lines(sub, [{"text": f"{topic_en}  ·  {topic_cn}" + (f"   |   {author}" if author else ""),
                       "size": 22, "color": GRAY, "align": PP_ALIGN.CENTER}],
                align=PP_ALIGN.CENTER)

    return prs


# ---------------------------------------------------------------- 入口
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--json", required=True)
    ap.add_argument("--out", required=True)
    ap.add_argument("--template", default=DEFAULT_TEMPLATE)
    args = ap.parse_args()

    with open(args.json, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation(args.template)
    # 模板自带一张空白首页时清掉
    while len(prs.slides._sldIdLst):
        rid = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[0]

    build(data, prs)

    os.makedirs(os.path.dirname(os.path.abspath(args.out)), exist_ok=True)
    prs.save(args.out)

    readme = os.path.splitext(args.out)[0] + "_素材清单.txt"
    with open(readme, "w", encoding="utf-8") as f:
        f.write("素材插入清单\n")
        f.write("=" * 52 + "\n")
        f.write("本课件已使用 PowerPoint 原生【图片占位符 / 媒体占位符】。\n")
        f.write("在 PowerPoint 中打开后，直接点击页面上的占位图标，\n")
        f.write("即可弹出【本地文件选择框】选择你电脑里的图片 / 音频 / 视频。\n")
        f.write("未插入素材的占位符，在放映和打印时会自动隐藏。\n")
        f.write("-" * 52 + "\n")
        for n in ASSET_NOTES:
            f.write(n + "\n")

    print(f"OK  slides={len(prs.slides)}  -> {args.out}")
    print(f"OK  assets readme -> {readme}")


if __name__ == "__main__":
    main()
