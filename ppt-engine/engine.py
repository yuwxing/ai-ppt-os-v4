import os
import random
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from pptx.enum.shapes import MSO_SHAPE

ANIMATIONS = ["fade", "fly_from_bottom", "zoom", "float", "split"]
TRANSITIONS = ["morph", "fade_smooth", "push", "reveal"]

LAYOUT_BUILDERS = {
    "cover": lambda slide, data, pptx: _build_cover(slide, data, pptx),
    "section": lambda slide, data, pptx: _build_section(slide, data, pptx),
    "content_text": lambda slide, data, pptx: _build_content_text(slide, data, pptx),
    "content_image": lambda slide, data, pptx: _build_content_image(slide, data, pptx),
    "content_split": lambda slide, data, pptx: _build_content_split(slide, data, pptx),
    "summary": lambda slide, data, pptx: _build_summary(slide, data, pptx),
}


def hex_to_rgb(hex_color: str):
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=(50, 50, 50), alignment=PP_ALIGN.LEFT,
                  font_name="微软雅黑"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _build_cover(slide, data, pptx):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    style = data.get("style", {})
    colors = style.get("color_scheme", {})
    primary = hex_to_rgb(colors.get("primary", "#2B5C8F"))
    fill.fore_color.rgb = RGBColor(*primary)

    title = data.get("title", "PPT标题")
    _add_text_box(slide, 1.5, 2.5, 9, 2, title, 40, True, (255, 255, 255), PP_ALIGN.CENTER)
    _add_text_box(slide, 1.5, 4.5, 9, 1, data.get("narrative", ""),
                  18, False, (220, 220, 220), PP_ALIGN.CENTER)


def _build_section(slide, data, pptx):
    style = data.get("style", {})
    colors = style.get("color_scheme", {})
    accent = hex_to_rgb(colors.get("accent", "#E8751A"))
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(*accent)
    _add_text_box(slide, 1.5, 2.5, 9, 2, data.get("title", ""),
                  36, True, (255, 255, 255), PP_ALIGN.CENTER)


def _build_content_text(slide, data, pptx):
    _add_text_box(slide, 0.5, 0.3, 9.5, 1, data.get("title", ""),
                  28, True, (50, 50, 50), PP_ALIGN.LEFT)
    items = data.get("content", [])
    body_text = "\n".join(f"• {item}" for item in items)
    _add_text_box(slide, 0.5, 1.5, 9.5, 5, body_text,
                  18, False, (60, 60, 60), PP_ALIGN.LEFT)


def _build_content_image(slide, data, pptx):
    _add_text_box(slide, 0.5, 0.3, 5, 1, data.get("title", ""),
                  28, True, (50, 50, 50), PP_ALIGN.LEFT)
    items = data.get("content", [])
    body_text = "\n".join(f"• {item}" for item in items[:3])
    _add_text_box(slide, 0.5, 1.5, 5, 3, body_text, 16, False, (60, 60, 60))

    img_path = data.get("image", {}).get("file_path")
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), Inches(6), Inches(1.2), Inches(4.5), Inches(5))


def _build_content_split(slide, data, pptx):
    _add_text_box(slide, 0.3, 0.3, 10, 0.8, data.get("title", ""),
                  26, True, (50, 50, 50), PP_ALIGN.CENTER)

    items = data.get("content", [])
    mid = len(items) // 2
    left_items = items[:mid] if len(items) > 2 else items
    right_items = items[mid:] if len(items) > 2 else []

    if left_items:
        _add_text_box(slide, 0.3, 1.5, 5, 4,
                      "\n".join(f"• {item}" for item in left_items),
                      16, False, (60, 60, 60), PP_ALIGN.LEFT)
    if right_items:
        _add_text_box(slide, 5.3, 1.5, 5, 4,
                      "\n".join(f"• {item}" for item in right_items),
                      16, False, (60, 60, 60), PP_ALIGN.LEFT)


def _build_summary(slide, data, pptx):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(240, 248, 255)
    _add_text_box(slide, 1.5, 1.5, 9, 1.5, "💡 总结",
                  36, True, (50, 50, 50), PP_ALIGN.CENTER)
    items = data.get("content", [])
    body_text = "\n".join(f"✦ {item}" for item in items[:5])
    _add_text_box(slide, 1.5, 3.5, 9, 4, body_text,
                  20, False, (60, 60, 60), PP_ALIGN.CENTER)


def build_ppt(ppt_data: dict, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    slides_data = ppt_data.get("slides", [])
    style = ppt_data.get("style", {})

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        slide_data["style"] = style

        layout = slide_data.get("layout", "content_text")
        builder = LAYOUT_BUILDERS.get(layout, LAYOUT_BUILDERS["content_text"])
        builder(slide, slide_data, prs)

        anim = slide_data.get("animation", {})
        if anim:
            _add_transition(slide, anim)

    prs.save(output_path)
    return output_path


def _add_transition(slide, anim):
    try:
        from pptx.oxml.ns import qn
        transition = slide.element.makeelement(qn("p:transition"), {})
        transition.set("advTm", "5000")
        slide.element.append(transition)
    except Exception:
        pass


def load_and_build(json_path: str, output_path: str = None) -> str:
    with open(json_path, "r", encoding="utf-8") as f:
        ppt_data = json.load(f)
    if output_path is None:
        output_path = json_path.replace(".json", ".pptx")
    return build_ppt(ppt_data, output_path)
