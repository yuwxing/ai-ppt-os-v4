import os, json
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from core.config import settings

class ExportAgent:
    async def run(self, qa_result: dict, topic: str) -> tuple:
        output_dir = Path(settings.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        safe_topic = "".join(c for c in topic if c.isalnum() or c in " _-").strip()[:30]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        file_name = f"{safe_topic}_{timestamp}.pptx"
        file_path = output_dir / file_name

        ppt_data = qa_result.get("ppt_data", {})
        slides_data = ppt_data.get("slides", [])
        template_file = ppt_data.get("template_file")

        games = qa_result.get("games", [])
        homework = qa_result.get("homework", [])
        theme = qa_result.get("theme_elevation", {})
        teacher = qa_result.get("teacher_guide", [])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        def add_bg(slide, color=RGBColor(0xF5, 0xF0, 0xEB)):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def add_shape(slide, left, top, width, height, fill_color, line_color=None):
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
            if line_color:
                shape.line.color.rgb = line_color
            else:
                shape.line.fill.background()
            return shape

        def add_text(slide, left, top, width, height, text, size=18, bold=False, color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.LEFT, font_name='微软雅黑'):
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.font.name = font_name
            p.alignment = align
            return txBox

        def add_multi_text(slide, left, top, width, height, items, size=16, color=RGBColor(0x33, 0x33, 0x33), spacing=Pt(8)):
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(size)
                p.font.color.rgb = color
                p.font.name = '微软雅黑'
                p.space_after = spacing
            return txBox

        PRIMARY = RGBColor(0x2B, 0x5C, 0x8F)
        ACCENT = RGBColor(0xE8, 0x75, 0x1A)
        DARK = RGBColor(0x33, 0x33, 0x33)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        LIGHT_BG = RGBColor(0xF5, 0xF0, 0xEB)
        LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xF8)

        # ── Cover ──
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, PRIMARY)
        add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5), topic, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        meta = qa_result.get('meta', {})
        subtitle = f"{meta.get('subject','')} · {meta.get('grade','')} · {meta.get('book','')}"
        if meta.get('lesson_type'):
            subtitle += f" · {meta['lesson_type']}"
        if meta.get('lesson_period'):
            subtitle += f" · {meta['lesson_period']}"
        add_text(slide, Inches(1), Inches(3.5), Inches(11), Inches(1), subtitle, size=20, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6), "AI-Wego · AI教师团队备课", size=14, color=RGBColor(0x99, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

        # ── Learning Objectives ──
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide)
        add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
        add_text(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), "🎯 学习目标", size=32, bold=True, color=WHITE)
        objectives = ["知识目标：掌握本节课的核心知识点", "能力目标：能够运用所学知识解决实际问题", "素养目标：培养学科思维与学习品质"]
        add_multi_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), [f"● {o}" for o in objectives], size=20, color=DARK, spacing=Pt(14))

        # ── Import / Story ──
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide)
        add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), ACCENT)
        add_text(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), "🎬 情境导入", size=32, bold=True, color=WHITE)
        narrative = ""
        for s in slides_data[:1]:
            narrative = s.get("narrative", "")
        import_text = narrative or "思考：这个情境与今天的学习有什么联系？"
        add_text(slide, Inches(1), Inches(2), Inches(11), Inches(4), import_text, size=22, color=DARK)
        add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8), "💡 自由发言 · 2分钟", size=18, color=ACCENT)

        # ── Task Slides ──
        task_titles = ["🧩 任务一：观察现象", "🧩 任务二：提出猜想", "🧩 任务三：实验验证", "🧩 任务四：规律总结", "🧩 任务五：生活应用"]
        task_activities = ["观察现象，记录数据", "基于观察提出假设", "设计实验验证猜想", "分析数据，得出结论", "将所学应用到生活实际"]
        for idx, task_title in enumerate(task_titles):
            slide = prs.slides.add_slide(blank_layout)
            add_bg(slide)
            colors = [RGBColor(0x4A, 0x90, 0xD9), RGBColor(0x5B, 0xB3, 0x8A), RGBColor(0xE8, 0x75, 0x1A), RGBColor(0x8E, 0x5C, 0xAE), RGBColor(0xD4, 0x5B, 0x5B)]
            add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), colors[idx % len(colors)])
            add_text(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), task_title, size=32, bold=True, color=WHITE)
            if idx < len(slides_data):
                content = slides_data[idx].get("content", [])
                texts = [c if isinstance(c, str) else c.get("text", "") for c in content] if isinstance(content, list) else [str(content)]
                add_multi_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(3.5), texts, size=18, color=DARK, spacing=Pt(10))
            add_text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6), f"⏱ {task_activities[idx]}  |  形式：小组合作", size=16, color=RGBColor(0x88, 0x88, 0x88))

        # ── Games ──
        if games:
            for g in games:
                slide = prs.slides.add_slide(blank_layout)
                add_bg(slide)
                add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0xE8, 0xA8, 0x3A))
                add_text(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), f"🎮 {g.get('name','课堂活动')}", size=32, bold=True, color=WHITE)
                add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6), f"类型：{g.get('type','')}  |  环节：{g.get('phase','')}  |  时长：{g.get('duration','')}", size=18, color=ACCENT)
                add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(4), g.get('description', ''), size=18, color=DARK)

        # ── Theme Elevation ──
        if theme:
            slide = prs.slides.add_slide(blank_layout)
            add_bg(slide, RGBColor(0xFD, 0xF0, 0xF0))
            add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0xD4, 0x5B, 0x8A))
            add_text(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), "🌈 主题升华", size=32, bold=True, color=WHITE)
            add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(0.6), f"核心价值：{theme.get('core_value','')}", size=18, bold=True, color=RGBColor(0xD4, 0x5B, 0x8A))
            add_text(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.5), f"形式：{theme.get('format','')}  |  时长：{theme.get('duration','')}", size=16, color=RGBColor(0xAA, 0x66, 0x88))
            add_text(slide, Inches(1), Inches(3.5), Inches(11), Inches(3.5), theme.get('content', ''), size=18, color=DARK)

        # ── Homework ──
        if homework:
            slide = prs.slides.add_slide(blank_layout)
            add_bg(slide)
            add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x3B, 0x8E, 0x5C))
            add_text(slide, Inches(0.8), Inches(0.2), Inches(11), Inches(0.8), "📋 课后作业", size=32, bold=True, color=WHITE)
            y = 1.8
            for tier in ['基础', '拓展', '实践']:
                items = [h for h in homework if h.get('tier') == tier]
                if not items: continue
                add_text(slide, Inches(0.8), Inches(y), Inches(11), Inches(0.5), f"【{tier}作业】", size=18, bold=True, color=PRIMARY)
                y += 0.5
                for h in items:
                    add_text(slide, Inches(1.2), Inches(y), Inches(11), Inches(0.4), f"• {h.get('title','')}（{h.get('estimated_time','')} · {h.get('difficulty','')}）", size=16, color=DARK)
                    y += 0.4

        # ── Thank You ──
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide, PRIMARY)
        add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5), "感谢聆听", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(1), "AI-Wego · AI教师团队", size=22, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

        prs.save(str(file_path))
        return str(file_path), file_name
